"""Before/After: 1 cerinta proasta → 3 cerinte IEEE 29148."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_BG    = RGBColor(0x1E, 0x1E, 0x2E)
DARK_CARD  = RGBColor(0x2A, 0x2B, 0x3D)
DARK_CARD2 = RGBColor(0x1A, 0x1B, 0x2A)
ACCENT     = RGBColor(0x89, 0xB4, 0xFA)
GREEN      = RGBColor(0xA6, 0xE3, 0xA1)
RED        = RGBColor(0xF3, 0x8B, 0xA8)
YELLOW     = RGBColor(0xF9, 0xE2, 0xAF)
ORANGE     = RGBColor(0xFA, 0xB3, 0x87)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xBA, 0xC2, 0xDE)
PURPLE     = RGBColor(0xCB, 0xA6, 0xF7)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide, color=DARK_BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def tb(slide, x, y, w, h, lines, fill=None):
    """lines = list of (text, size, bold, color, align)"""
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    first = True
    for item in lines:
        text, size, bold, color = item[0], item[1], item[2], item[3]
        align = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return shape


def arrow(slide, x1, y1, x2, y2, color=ACCENT):
    """Simple line connector approximated with a thin textbox."""
    # Using a shape connector
    from pptx.util import Inches
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)


# =========================================================
# SLIDE 1 — Title
# =========================================================
sl = prs.slides.add_slide(BLANK)
bg(sl)

tb(sl, 1.5, 2.0, 10, 1.0,
   [("Before / After", 48, True, ACCENT, PP_ALIGN.CENTER)])
tb(sl, 1.5, 3.1, 10, 0.7,
   [("1 cerință proastă  →  3 cerințe IEEE 29148", 26, False, WHITE, PP_ALIGN.CENTER)])
tb(sl, 1.5, 4.1, 10, 0.5,
   [("Cum arată diferența concret, pe un exemplu real", 16, False, GRAY, PP_ALIGN.CENTER)])

# decorative bar
bar = sl.shapes.add_textbox(Inches(4.5), Inches(3.9), Inches(4.33), Inches(0.04))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT


# =========================================================
# SLIDE 2 — BEFORE: cerinta proasta
# =========================================================
sl = prs.slides.add_slide(BLANK)
bg(sl)

tb(sl, 0.3, 0.1, 8, 0.45,
   [("BEFORE — 1 cerință (prea mare, prea vagă, câmpuri lipsă)", 17, True, RED)])
tb(sl, 8.5, 0.1, 4.5, 0.45,
   [("Problemele identificate", 17, True, RED)])

# the bad requirement
bad_req = """\
---
id: AUTH-001
status: confirmed
owner: Alex
---

# User Authentication

> WHY: users need to log in and manage their accounts.

## WHAT — Contract
- The system shall allow users to log in with
  email and password, reset their password via
  email link, and automatically log out after
  30 minutes of inactivity.
- It shall also rate-limit failed attempts and
  lock the account after 5 failures.
- Error messages shall be user-friendly and not
  reveal security details.

## HOW — Acceptance (= tests)
AC-1  valid login → token returned
AC-2  wrong password → error shown
AC-3  password reset email sent
AC-4  reset link expires after 24h
AC-5  session expires after 30min inactivity
AC-6  5 failed attempts → account locked
AC-7  locked account → informative message shown
AC-8  rate limiter blocks >10 req/min per IP
AC-9  error message doesn't reveal if email exists
AC-10 locked account can be unlocked by admin
"""

tb(sl, 0.3, 0.62, 7.8, 6.7,
   [(bad_req, 9.5, False, RGBColor(0xFF, 0xCC, 0xCC))],
   fill=DARK_CARD2)

# problem annotations
problems = [
    (RED,    "✗  Singular VIOLAT",
             "'shall allow ... log in AND reset ... AND log out' = 3 capabilități diferite într-un singur bullet"),
    (RED,    "✗  10 AC (max recomandat: 7)",
             "Semnalul clar că cerința acoperă prea mult — ar trebui ≥3 cerințe separate"),
    (RED,    "✗  Source lipsă",
             "Nu știm cine a cerut asta (Product Owner? Legal? Security team?)"),
    (RED,    "✗  Priority lipsă",
             "Login și rate-limiting au priorități diferite — nu se pot planifica separat"),
    (RED,    "✗  Feasibility lipsă",
             "Nu există estimare de complexitate / risc tehnic"),
    (RED,    "✗  Rationale per statement lipsă",
             "De ce 30 min? De ce 5 încercări? Nu există justificare per regulă"),
    (RED,    "✗  Verifiable method lipsă",
             "AC-7 și AC-9 se testează manual sau automat? Nu e declarat"),
    (YELLOW, "⚠  Unambiguous parțial",
             "'user-friendly' este vag — nu există criteriu testabil"),
]

y_p = 0.62
for bg_c, title, desc in problems:
    tb(sl, 8.5, y_p, 4.6, 0.7,
       [(title, 10, True, bg_c),
        (desc,  9, False, GRAY)],
       fill=DARK_CARD)
    y_p += 0.75


# =========================================================
# SLIDE 3 — THE SPLIT visual
# =========================================================
sl = prs.slides.add_slide(BLANK)
bg(sl)

tb(sl, 0.3, 0.1, 12.5, 0.5,
   [("Cum se împarte: 1 → 3 cerințe atomice", 22, True, ACCENT, PP_ALIGN.CENTER)])

# left: the 1
tb(sl, 0.4, 1.2, 3.6, 4.5,
   [("AUTH-001\n(original)", 15, True, RED, PP_ALIGN.CENTER),
    ("\n• Login cu email/parola\n• Reset parola via email\n• Logout automat 30min\n• Rate-limit attempts\n• Lock dupa 5 fail-uri\n• Error messages\n\n10 AC\n0 câmpuri standard", 11, False, GRAY)],
   fill=DARK_CARD2)

# arrow area label
tb(sl, 4.15, 3.1, 1.0, 0.5,
   [("split\n→", 18, True, YELLOW, PP_ALIGN.CENTER)])

# right: 3 boxes
splits = [
    (ACCENT,  "AUTH-LOGIN-001",
              "User Login",
              "Capabilitate: autentificare cu\nemail + parola\n\n• source: Product Owner\n• priority: must\n• feasibility: high\n\n3 AC (login ok, login fail,\nenum protection)"),
    (GREEN,   "AUTH-RESET-002",
              "Password Reset",
              "Capabilitate: resetare parola\nvia link email\n\n• source: Product Owner\n• priority: must\n• feasibility: medium\n\n3 AC (link trimis, link expirat,\nparola schimbata)"),
    (PURPLE,  "AUTH-SESSION-003",
              "Session & Rate Limiting",
              "Capabilitate: expirare sesiune\n+ rate-limit + lock cont\n\n• source: Security Team\n• priority: should\n• feasibility: high\n\n4 AC (timeout, rate-limit,\nlock, admin unlock)"),
]

y_boxes = [0.9, 2.9, 4.9]
for i, ((color, req_id, req_title, req_body), y_b) in enumerate(zip(splits, y_boxes)):
    tb(sl, 5.4, y_b, 7.5, 1.75,
       [(f"{req_id} — {req_title}", 13, True, color),
        (req_body, 10, False, GRAY)],
       fill=DARK_CARD)

tb(sl, 0.3, 6.9, 12.5, 0.45,
   [("Fiecare cerință are o singură capabilitate, câmpuri complete și ≤7 AC — lint trece fără avertizări", 13, False, GRAY, PP_ALIGN.CENTER)])


# =========================================================
# SLIDE 4 — AFTER: AUTH-LOGIN-001
# =========================================================
sl = prs.slides.add_slide(BLANK)
bg(sl)

tb(sl, 0.3, 0.08, 6, 0.42,
   [("AFTER — Cerința 1/3:  AUTH-LOGIN-001", 17, True, ACCENT)])
tb(sl, 6.5, 0.08, 6.5, 0.42,
   [("Ce s-a adăugat față de 'before'", 17, True, GREEN)])

req1 = """\
---
id:          AUTH-LOGIN-001
title:       User login with email/password
status:      confirmed
layer:       feature
owner:       Alex
source:      Product Owner — sprint 2026-05-01
priority:    must
feasibility: high
depends_on:  [CORE-SESSION-002]
---

# User login

> WHY: users need to reach their own data securely.
> SOURCE RATIONALE: login is the primary entry point
>   per product spec v1.2, section 3.

## WHAT — Contract (normative)
- [S1] The system shall authenticate a registered user
  via email address and password and return a signed
  session token within 2 seconds under normal load.
  Rationale: 2s SLA from UX research — abandonment
  spikes above this threshold.

- [S2] The system shall reject an unregistered email
  with a generic HTTP 401 error without revealing
  whether the email address exists in the system.
  Rationale: prevents user-enumeration (OWASP A2).

## HOW — Acceptance
AC-1  [verifiable by: automated test]
  Given  registered user, correct credentials
  When   POST /auth/login
  Then   HTTP 200, signed JWT, latency < 2s

AC-2  [verifiable by: automated test]
  Given  wrong password
  When   POST /auth/login
  Then   HTTP 401, generic error body

AC-3  [verifiable by: security test]
  Given  attacker with known/unknown emails
  When   50 login attempts, mixed valid/invalid
  Then   response time and body are identical
"""

tb(sl, 0.3, 0.58, 6.2, 6.75,
   [(req1, 9.3, False, GREEN)],
   fill=DARK_CARD2)

# what's new annotations
news1 = [
    (GREEN,  "✓  title: câmp nou",          "Titlu lizibil separat de id — util în rapoarte și traceability"),
    (GREEN,  "✓  source: câmp nou",          "Product Owner, sprint, dată — știm exact cine a cerut-o"),
    (GREEN,  "✓  priority: must",            "Planificabil: login e 'must', alte cerinte pot fi 'should'"),
    (GREEN,  "✓  feasibility: high",         "Echipa de dev știe că nu există risc tehnic major"),
    (GREEN,  "✓  SOURCE RATIONALE:",         "De ce există cerința — legătură la spec extern"),
    (GREEN,  "✓  Rationale: per statement",  "S1: de ce 2s? S2: de ce generic error? Justificat inline"),
    (GREEN,  "✓  [S1], [S2] sub-ids",        "Fiecare statement are un id propriu pentru traceability"),
    (GREEN,  "✓  [verifiable by: ...]",      "Metoda de verificare declarată pe fiecare AC"),
    (YELLOW, "◈  Singular respectat",        "O singură capabilitate: login. Reset și session = cerinte separate"),
    (YELLOW, "◈  3 AC (era 10 în 'before')", "Compact, testabil, fără overlap cu celelalte cerinte"),
]

y_n = 0.58
for color, title, desc in news1:
    tb(sl, 6.6, y_n, 6.5, 0.6,
       [(title, 10, True, color),
        (desc,  9,  False, GRAY)],
       fill=DARK_CARD)
    y_n += 0.65


# =========================================================
# SLIDE 5 — AFTER: AUTH-RESET-002
# =========================================================
sl = prs.slides.add_slide(BLANK)
bg(sl)

tb(sl, 0.3, 0.08, 6, 0.42,
   [("AFTER — Cerința 2/3:  AUTH-RESET-002", 17, True, GREEN)])
tb(sl, 6.5, 0.08, 6.5, 0.42,
   [("Observații cheie", 17, True, GREEN)])

req2 = """\
---
id:          AUTH-RESET-002
title:       Password reset via email link
status:      confirmed
layer:       feature
owner:       Alex
source:      Product Owner — sprint 2026-05-01
priority:    must
feasibility: medium
depends_on:  [AUTH-LOGIN-001, CORE-EMAIL-005]
---

# Password reset

> WHY: users who forget their password need a
>   secure, self-service recovery path.
> SOURCE RATIONALE: reduces support tickets —
>   currently 30% of support volume per ops report.

## WHAT — Assumptions & constraints
- Assumes SMTP delivery is reliable (CORE-EMAIL-005).
- Out of scope: social login / OAuth recovery.

## WHAT — Contract (normative)
- [S1] The system shall send a password-reset link
  to the registered email address within 60 seconds
  of a valid reset request.
  Rationale: 60s balances SMTP latency with UX
  expectations (user research Q1-2026).

- [S2] The system shall invalidate the reset link
  after 24 hours or after first use, whichever
  comes first.
  Rationale: limits attack window per OWASP A7.

## HOW — Acceptance
AC-1  [verifiable by: automated test]
  Given  registered email
  When   POST /auth/reset-request
  Then   email received within 60s, HTTP 202

AC-2  [verifiable by: automated test]
  Given  a valid reset link
  When   used once, then used again
  Then   second use returns HTTP 410 Gone

AC-3  [verifiable by: automated test]
  Given  a reset link older than 24 hours
  When   accessed
  Then   HTTP 410 Gone, no password change
"""

tb(sl, 0.3, 0.58, 6.2, 6.75,
   [(req2, 9.3, False, GREEN)],
   fill=DARK_CARD2)

obs2 = [
    (GREEN,  "✓  Assumptions & constraints",
             "SMTP dependency declarată explicit — Socrate ar fi mulțumit"),
    (GREEN,  "✓  'Out of scope' declarat",
             "OAuth recovery exclusă explicit — previne scope creep"),
    (GREEN,  "✓  source + priority separate",
             "priority: must — la fel ca login; planificabil în același sprint"),
    (GREEN,  "✓  feasibility: medium",
             "SMTP integration adaugă risc — echipa știe din start"),
    (GREEN,  "✓  Rationale: per statement",
             "S1: de ce 60s? S2: de ce 24h? Fiecare număr magic are justificare"),
    (GREEN,  "✓  depends_on: explicit",
             "Depinde de AUTH-LOGIN-001 și CORE-EMAIL-005 — graf clar"),
    (YELLOW, "◈  3 AC clare, fără overlap",
             "Fiecare AC testează un singur behavior; nu există AC duplicate cu login"),
    (YELLOW, "◈  Singular respectat",
             "Doar reset. Login și session management = cerinte separate"),
]

y_n = 0.58
for color, title, desc in obs2:
    tb(sl, 6.6, y_n, 6.5, 0.72,
       [(title, 10, True, color),
        (desc,  9,  False, GRAY)],
       fill=DARK_CARD)
    y_n += 0.77


# =========================================================
# SLIDE 6 — AFTER: AUTH-SESSION-003
# =========================================================
sl = prs.slides.add_slide(BLANK)
bg(sl)

tb(sl, 0.3, 0.08, 6, 0.42,
   [("AFTER — Cerința 3/3:  AUTH-SESSION-003", 17, True, PURPLE)])
tb(sl, 6.5, 0.08, 6.5, 0.42,
   [("Observații cheie", 17, True, PURPLE)])

req3 = """\
---
id:          AUTH-SESSION-003
title:       Session expiry and brute-force protection
status:      confirmed
layer:       feature
owner:       Alex
source:      Security Team — threat model review 2026-04
priority:    should
feasibility: high
depends_on:  [AUTH-LOGIN-001]
---

# Session & rate limiting

> WHY: protect accounts from brute-force and
>   session-hijacking attacks.
> SOURCE RATIONALE: security team threat model
>   identified session persistence and BF as
>   top-2 risks for the auth surface.

## WHAT — Assumptions & constraints
- Rate limiting operates at IP + account level.
- Admin unlock is out of scope for v1 (manual DB).

## WHAT — Contract (normative)
- [S1] The system shall invalidate a session token
  after 30 minutes of inactivity.
  Rationale: NIST SP 800-63B recommends ≤30min
  idle timeout for moderate assurance level.

- [S2] The system shall block login attempts from
  an IP address for 15 minutes after 10 failed
  attempts within any 5-minute window.
  Rationale: limits BF to <100 attempts/hour
  per IP (threat model TM-AUTH-02).

- [S3] The system shall lock an account after
  5 consecutive failed login attempts and require
  admin action to unlock.
  Rationale: account-level lock prevents
  distributed BF across many IPs.

## HOW — Acceptance
AC-1  [verifiable by: automated test]
  Given  authenticated session
  When   no activity for 31 minutes
  Then   next request returns HTTP 401

AC-2  [verifiable by: load test]
  Given  single IP
  When   11 failed logins within 5 minutes
  Then   IP blocked for 15min, HTTP 429

AC-3  [verifiable by: automated test]
  Given  an account
  When   5 consecutive wrong passwords
  Then   account locked, subsequent attempts
         return HTTP 423 regardless of password

AC-4  [verifiable by: manual test]
  Given  a locked account
  When   admin calls DELETE /admin/lock/{userId}
  Then   account unlocked, login works again
"""

tb(sl, 0.3, 0.58, 6.2, 6.75,
   [(req3, 9.3, False, PURPLE)],
   fill=DARK_CARD2)

obs3 = [
    (PURPLE, "✓  source: Security Team",
             "Sursa diferita fata de login/reset (Product Owner) — prioritate si context diferit"),
    (PURPLE, "✓  priority: should (nu must)",
             "Rate-limiting e important dar nu blocant pentru MVP — planificabil separat"),
    (GREEN,  "✓  Rationale cu referinte",
             "S1: NIST SP 800-63B citat explicit. S2: threat model TM-AUTH-02"),
    (GREEN,  "✓  Assumptions: Admin unlock v1",
             "'manual DB' declarat explicit — nu o surpriza pentru tester"),
    (GREEN,  "✓  4 AC, fiecare cu metoda",
             "AC-2: load test (nu unit test!). AC-4: manual test — metode distincte"),
    (GREEN,  "✓  Singular: 3 statements legate",
             "Toate privesc protectia sesiunii/BF — o singura capabilitate coerenta"),
    (YELLOW, "◈  depends_on: AUTH-LOGIN-001",
             "Nu poate exista fara login — graful de dependente e corect si complet"),
    (YELLOW, "◈  'Admin unlock out of scope v1'",
             "Declarat explicit, nu lasat vag. Urmatoarea versiune stie exact ce sa faca"),
]

y_n = 0.58
for color, title, desc in obs3:
    tb(sl, 6.6, y_n, 6.5, 0.72,
       [(title, 10, True, color),
        (desc,  9,  False, GRAY)],
       fill=DARK_CARD)
    y_n += 0.77


# =========================================================
# SLIDE 7 — Summary comparison table
# =========================================================
sl = prs.slides.add_slide(BLANK)
bg(sl)

tb(sl, 0.3, 0.08, 12.5, 0.45,
   [("Comparație directă: BEFORE vs AFTER", 22, True, ACCENT, PP_ALIGN.CENTER)])

headers = ["Criteriu", "BEFORE\nAUTH-001", "AFTER\nAUTH-LOGIN-001", "AFTER\nAUTH-RESET-002", "AFTER\nAUTH-SESSION-003"]
col_x   = [0.3, 3.55, 5.8, 8.05, 10.3]
col_w   = [3.2, 2.2,  2.2, 2.2,  2.8]

# header
for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    fill = RGBColor(0x50, 0x50, 0x80) if i == 1 else DARK_CARD
    color = RED if i == 1 else ACCENT
    tb(sl, x, 0.6, w, 0.55,
       [(h, 10, True, color, PP_ALIGN.CENTER)],
       fill=fill)

rows = [
    ("Singular (1 capabilitate)",         "✗  3 în 1",        "✓  login only",      "✓  reset only",      "✓  session/BF only"),
    ("Număr AC",                          "✗  10 AC",         "✓  3 AC",            "✓  3 AC",            "✓  4 AC"),
    ("source: (cine a cerut)",            "✗  lipsă",         "✓  Product Owner",   "✓  Product Owner",   "✓  Security Team"),
    ("priority:",                         "✗  lipsă",         "✓  must",            "✓  must",            "✓  should"),
    ("feasibility:",                      "✗  lipsă",         "✓  high",            "✓  medium",          "✓  high"),
    ("Rationale per statement",           "✗  lipsă",         "✓  S1, S2",          "✓  S1, S2",          "✓  S1, S2, S3"),
    ("[verifiable by:] per AC",           "✗  lipsă",         "✓  toate 3",         "✓  toate 3",         "✓  toate 4"),
    ("Assumptions & constraints",         "✗  lipsă",         "—  n/a",             "✓  SMTP, scope",     "✓  IP+acct, v1"),
    ("depends_on explicit",               "✗  lipsă",         "✓  CORE-SESSION",    "✓  LOGIN + EMAIL",   "✓  LOGIN"),
    ("Număr cerinte IEEE 29148 acoperite","~4/13  (31%)",      "~11/13  (85%)",      "~12/13  (92%)",      "~12/13  (92%)"),
]

y_r = 1.22
for row in rows:
    is_last = row == rows[-1]
    for i, (val, x, w) in enumerate(zip(row, col_x, col_w)):
        if i == 0:
            color = WHITE
            fill  = DARK_CARD
        elif i == 1:
            color = RED   if "✗" in val else GRAY
            fill  = RGBColor(0x35, 0x20, 0x28)
        else:
            color = GREEN if "✓" in val else (GRAY if "—" in val else YELLOW)
            fill  = DARK_CARD
        if is_last:
            color = YELLOW if i == 1 else GREEN
            fill  = RGBColor(0x28, 0x30, 0x28)
        tb(sl, x, y_r, w, 0.4 if not is_last else 0.52,
           [(val, 10 if not is_last else 11, is_last, color, PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)],
           fill=fill)
    y_r += 0.42 if not is_last else 0.54

import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "before_after_ieee29148.pptx")
prs.save(out)
print(f"saved: {out}")
