"""Generate IEEE 29148 comparison presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ---- palette ----
DARK_BG    = RGBColor(0x1E, 0x1E, 0x2E)   # dark slate
ACCENT     = RGBColor(0x89, 0xB4, 0xFA)   # blue
GREEN      = RGBColor(0xA6, 0xE3, 0xA1)   # green
RED        = RGBColor(0xF3, 0x8B, 0xA8)   # red
YELLOW     = RGBColor(0xF9, 0xE2, 0xAF)   # yellow
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xBA, 0xC2, 0xDE)
DARK_CARD  = RGBColor(0x31, 0x32, 0x44)
DARK_CARD2 = RGBColor(0x28, 0x29, 0x3B)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank


def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill_color=DARK_CARD, text="", font_size=14,
        bold=False, font_color=WHITE, align=PP_ALIGN.LEFT, wrap=True, radius=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = font_color
    # background fill on the shape
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = fill_color
    return txBox


def label(slide, x, y, w, h, text, size=14, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, fill=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if fill:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = fill
    else:
        txBox.fill.background()
    return txBox


def multiline_box(slide, x, y, w, h, lines, fill_color=DARK_CARD):
    """lines = list of (text, size, bold, color)"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = fill_color
    first = True
    for (text, size, bold, color) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox


# ============================================================
# SLIDE 1 — Title
# ============================================================
sl = prs.slides.add_slide(BLANK)
bg(sl)

label(sl, 1.5, 1.8, 10, 1.2, "Cum arată un requirement", 42, True, ACCENT, PP_ALIGN.CENTER)
label(sl, 1.5, 2.9, 10, 0.8, "dacă respectă ≥80% din IEEE 29148", 28, False, WHITE, PP_ALIGN.CENTER)
label(sl, 1.5, 3.8, 10, 0.5, "Comparație: format actual  →  format standard", 18, False, GRAY, PP_ALIGN.CENTER)

# decorative line
line_box = sl.shapes.add_textbox(Inches(4), Inches(3.55), Inches(5.33), Pt(2))
line_box.fill.solid()
line_box.fill.fore_color.rgb = ACCENT

label(sl, 1.5, 5.5, 10, 0.6,
      "IEEE 29148:2018 — Systems and Software Engineering — Life cycle processes — Requirements engineering",
      12, False, GRAY, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 — IEEE 29148: cele 9 caracteristici individuale
# ============================================================
sl = prs.slides.add_slide(BLANK)
bg(sl)

label(sl, 0.4, 0.2, 12.5, 0.6, "IEEE 29148 — Cele 9 caracteristici ale unui requirement individual", 22, True, ACCENT)

chars = [
    ("1  Necessary",    "Lipsa lui creează un deficit real. Nu există cerințe speculative."),
    ("2  Appropriate",  "Nivelul de abstractizare potrivit audienței (stakeholder vs. developer)."),
    ("3  Unambiguous",  "O singură interpretare posibilă. Niciun termen vag fără definiție."),
    ("4  Complete",     "Declarat complet — nu necesită informații externe pentru înțelegere."),
    ("5  Singular",     "O singură capabilitate sau condiție per cerință. Nu 'și A și B'."),
    ("6  Feasible",     "Realizabil tehnic și economic în constrângerile proiectului."),
    ("7  Verifiable",   "Poate fi verificat prin test, inspecție, analiză sau demonstrație."),
    ("8  Correct",      "Reflectă fidel nevoia stakeholderului, fără distorsionare."),
    ("9  Conforming",   "Respectă un template și o convenție de stilizare (ex: 'shall')."),
]

cols = 3
rows = 3
card_w = 3.9
card_h = 1.5
x_start = 0.35
y_start = 1.0
gap_x = 0.15
gap_y = 0.18

for i, (title, desc) in enumerate(chars):
    col = i % cols
    row = i // cols
    x = x_start + col * (card_w + gap_x)
    y = y_start + row * (card_h + gap_y)
    multiline_box(sl, x, y, card_w, card_h, [
        (title, 13, True,  YELLOW),
        (desc,  11, False, GRAY),
    ], fill_color=DARK_CARD)


# ============================================================
# SLIDE 3 — Format ACTUAL (requirement-manager azi)
# ============================================================
sl = prs.slides.add_slide(BLANK)
bg(sl)

label(sl, 0.3, 0.15, 6, 0.5, "Format ACTUAL — requirement-manager", 20, True, ACCENT)
label(sl, 6.7, 0.15, 6, 0.5, "Ce caracteristici acoperă?", 20, True, ACCENT)

# left: code-like example
code = """\
---
id: AUTH-LOGIN-001
status: confirmed
layer: feature
owner: Alex
depends_on: [CORE-SESSION-002]
---

# User login

> WHY: users need to reach their data securely.

## WHAT — Contract (normative)
- It shall accept email + password and return
  a session token.
- It shall reject unknown email with a generic
  error (no user enumeration).

## HOW — Acceptance (= tests)
AC-1
  Given  a registered user
  When   correct credentials are submitted
  Then   a session token is returned

AC-2
  Given  an unknown email
  When   login is attempted
  Then   a generic error is returned, HTTP 401
"""

multiline_box(sl, 0.3, 0.75, 6.1, 6.5,
    [(code, 10, False, GREEN)],
    fill_color=RGBColor(0x18, 0x18, 0x28))

# right: checklist
checks = [
    (True,  "Necessary     — WHY section"),
    (True,  "Unambiguous   — 'shall' phrasing ghidată"),
    (True,  "Singular      — un bullet = o regulă"),
    (True,  "Verifiable    — AC Given/When/Then"),
    (True,  "Conforming    — template impus"),
    (True,  "Complete (parțial) — contract normativ"),
    (False, "Appropriate   — lipsă nivel de abstractizare explicit"),
    (False, "Feasible      — nu există câmp"),
    (False, "Correct       — nu există sursă/stakeholder"),
    (False, "Priority      — lipsă (MoSCoW / numeric)"),
    (False, "Source        — lipsă (cine a cerut-o)"),
    (False, "Rationale per statement — WHY e global, nu per bullet"),
    (False, "Version/Baseline — lipsă"),
]

y_c = 0.75
for ok, text in checks:
    color = GREEN if ok else RED
    symbol = "✓" if ok else "✗"
    multiline_box(sl, 6.7, y_c, 6.3, 0.38,
        [(f"{symbol}  {text}", 11, False, color)],
        fill_color=DARK_CARD2)
    y_c += 0.42

covered = sum(1 for ok, _ in checks if ok)
label(sl, 6.7, y_c + 0.05, 6.3, 0.4,
      f"Acoperire: {covered}/{len(checks)}  ≈ {int(covered/len(checks)*100)}%",
      13, True, YELLOW, PP_ALIGN.CENTER, fill=DARK_CARD)


# ============================================================
# SLIDE 4 — Format IEEE 29148 ≥80% (cum ar arăta)
# ============================================================
sl = prs.slides.add_slide(BLANK)
bg(sl)

label(sl, 0.3, 0.15, 12.5, 0.5, "Format IEEE 29148 ≥80% — cum ar arăta requirement-ul extins", 20, True, ACCENT)

code2 = """\
---
id:           AUTH-LOGIN-001
title:        User login with email/password
status:       confirmed
layer:        feature
owner:        Alex
source:       Product Owner — sprint planning 2026-05-01    # ← NOU: cine a cerut-o
priority:     must                                          # ← NOU: MoSCoW
feasibility:  high                                         # ← NOU: realizabil tehnic
depends_on:   [CORE-SESSION-002]
---

# User login

> WHY: users need to reach their own data securely.
> SOURCE RATIONALE: legal requires session-token auth
>   per compliance spec v2.3, section 4.1.           # ← NOU: rationale legat de sursă

## WHAT — Contract (normative)
- [AUTH-LOGIN-001-S1] The system shall accept a registered
  user's email address and password and return a signed
  session token within 2 seconds.
  Rationale: 2s SLA derived from UX research —         # ← NOU: rationale per statement
  abandonment rate spikes above 2s.

- [AUTH-LOGIN-001-S2] The system shall reject an unregistered
  email with a generic error message (HTTP 401) without
  revealing whether the email is registered.
  Rationale: prevents user enumeration attacks (OWASP A2).

## WHAT — Assumptions & constraints            # ← NOU: secțiune IEEE 29148 §9.5
- Assumes SMTP delivery is available for password reset.
- Constrained to password auth only (OAuth deferred to v2).

## HOW — Acceptance (= tests)
AC-1  [verifiable by: automated test]          # ← NOU: metoda de verificare
  Given  a registered user with correct credentials
  When   POST /auth/login is called
  Then   HTTP 200, body contains signed JWT, latency < 2s

AC-2  [verifiable by: automated test]
  Given  an email not in the user table
  When   POST /auth/login is called
  Then   HTTP 401, body = generic error, no user-specific info

AC-3  [verifiable by: penetration test]
  Given  an attacker enumerating emails
  When   100 login attempts with known/unknown emails
  Then   response time and body are indistinguishable

## Traceability                                # ← NOU: trasabilitate upstream
- Upstream:   stakeholder-need SN-004 (secure access)
- Downstream: CORE-SESSION-002 (auto-linked via depends_on)
"""

multiline_box(sl, 0.3, 0.75, 8.5, 6.5,
    [(code2, 9.5, False, GREEN)],
    fill_color=RGBColor(0x18, 0x18, 0x28))

# right: what's new
news = [
    ("NOU — câmpuri frontmatter", [
        "source: — cine a cerut cerința",
        "priority: must/should/could",
        "feasibility: high/medium/low",
    ]),
    ("NOU — rationale per statement", [
        "Fiecare bullet 'shall' are un 'Rationale:' inline",
        "Explică DE CE acel comportament specific",
    ]),
    ("NOU — Assumptions & constraints", [
        "Secțiune dedicată pentru presupuneri",
        "și limite de scop ale cerinței",
    ]),
    ("NOU — metoda de verificare pe AC", [
        "[verifiable by: automated test]",
        "[verifiable by: penetration test]",
        "[verifiable by: inspection]",
    ]),
    ("NOU — Traceability upstream", [
        "Link la nevoia stakeholderului (SN-xxx)",
        "Nu doar downstream la cod",
    ]),
]

y_n = 0.75
for title, items in news:
    lines = [(title, 12, True, YELLOW)]
    for it in items:
        lines.append(("  • " + it, 10, False, GRAY))
    h = 0.3 + len(items) * 0.28
    multiline_box(sl, 9.1, y_n, 4.0, h, lines, fill_color=DARK_CARD)
    y_n += h + 0.12


# ============================================================
# SLIDE 5 — Gap analysis table
# ============================================================
sl = prs.slides.add_slide(BLANK)
bg(sl)

label(sl, 0.3, 0.15, 12.5, 0.5, "Gap Analysis — ce adaugă IEEE 29148 față de formatul actual", 20, True, ACCENT)

headers = ["Caracteristică IEEE 29148", "Format actual", "Format ≥80%", "Efort adăugare"]
col_w   = [3.4, 2.6, 3.8, 2.1]
x_cols  = [0.3, 3.75, 6.4, 10.25]

# header row
for i, (h_text, w) in enumerate(zip(headers, col_w)):
    multiline_box(sl, x_cols[i], 0.75, w, 0.4,
        [(h_text, 12, True, DARK_BG)],
        fill_color=ACCENT)

rows_data = [
    ("Necessary (WHY)",              "✓ secțiunea WHY",          "✓ identic",                     "—"),
    ("Singular (un bullet = 1 req)", "✓ ghidat prin template",   "✓ + lint enforcement",           "mic"),
    ("Unambiguous ('shall')",        "✓ ghidat",                  "✓ + lint shall-check",           "mic"),
    ("Verifiable (AC G/W/T)",        "✓ AC Given/When/Then",     "✓ + [verifiable by: method]",    "mic"),
    ("Conforming (template)",        "✓ template fix",            "✓ identic",                      "—"),
    ("Complete (contract)",          "✓ parțial",                 "✓ + Assumptions & constraints",  "mediu"),
    ("Correct — Source/Stakeholder", "✗ lipsă",                   "✓ câmp source: în frontmatter",  "mic"),
    ("Priority (MoSCoW)",            "✗ lipsă",                   "✓ câmp priority: în frontmatter","mic"),
    ("Feasible",                     "✗ lipsă",                   "✓ câmp feasibility:",            "mic"),
    ("Rationale per statement",      "✗ doar WHY global",         "✓ Rationale: inline per bullet", "mediu"),
    ("Traceability upstream",        "✗ doar downstream (cod)",   "✓ secțiune Traceability",        "mediu"),
    ("Verif. method per AC",         "✗ lipsă",                   "✓ [verifiable by: x] pe AC",     "mic"),
    ("Assumptions & constraints",    "✗ lipsă",                   "✓ secțiune dedicată",            "mic"),
]

y_r = 1.22
for (char, actual, target, effort) in rows_data:
    vals = [char, actual, target, effort]
    colors = [WHITE, (GREEN if "✓" in actual else RED), GREEN,
              (GRAY if effort == "—" else (GREEN if effort == "mic" else YELLOW))]
    for i, (val, col_color) in enumerate(zip(vals, colors)):
        multiline_box(sl, x_cols[i], y_r, col_w[i], 0.36,
            [(val, 10, False, col_color)],
            fill_color=DARK_CARD if y_r % 0.8 < 0.4 else DARK_CARD2)
    y_r += 0.38


# ============================================================
# SLIDE 6 — Concluzie + recomandare
# ============================================================
sl = prs.slides.add_slide(BLANK)
bg(sl)

label(sl, 0.3, 0.15, 12.5, 0.6, "Concluzie și recomandare", 22, True, ACCENT)

summary = [
    ("Ce ai deja (≈46%)", GREEN, [
        "Template cu secțiuni WHY / WHAT / HOW / WHERE",
        "Phrasing 'shall' ghidat în comentarii",
        "AC în format Given/When/Then",
        "id, status, layer, owner, depends_on în frontmatter",
        "lint pentru lizibilitate și structură",
    ]),
    ("Ce lipsește pentru ≥80%", YELLOW, [
        "source: (cine a cerut cerința — stakeholder)",
        "priority: must / should / could / won't",
        "feasibility: high / medium / low",
        "Rationale per statement (nu doar WHY global)",
        "Traceability upstream (link la nevoia stakeholderului)",
        "[verifiable by: method] pe fiecare AC",
        "Secțiune Assumptions & constraints",
    ]),
    ("Efort estimat", ACCENT, [
        "6 câmpuri noi în frontmatter   → ~30 min (template + lint)",
        "Rationale inline per bullet    → ~1h (template + ghid)",
        "Traceability upstream          → ~2h (nou câmp + link checker)",
        "[verifiable by] pe AC          → ~1h (template + lint check)",
        "Total: ~4-5h implementare      → de la 46% la ≈83% conformanță",
    ]),
]

x_s = 0.3
for (title, title_color, items) in summary:
    lines = [(title, 14, True, title_color)]
    for it in items:
        lines.append(("  • " + it, 11, False, GRAY))
    h = 0.4 + len(items) * 0.35
    multiline_box(sl, x_s, 0.85, 4.1, h, lines, fill_color=DARK_CARD)
    x_s += 4.35

label(sl, 0.3, 6.1, 12.5, 0.5,
      "Senate verdict: MODIFY — revendicarea 'IEEE 29148-aligned' este prematură la 46%. "
      "La ≥80% (după ~5h efort) devine legitimă cu un disclaimer explicit pe ce se verifică.",
      13, False, GRAY, PP_ALIGN.CENTER)

# ============================================================
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ieee29148_comparison.pptx")
prs.save(out)
print(f"saved: {out}")
